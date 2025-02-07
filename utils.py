# utils.py
import os
import pandas as pd
from bs4 import BeautifulSoup
import docx
import PyPDF2
from pptx import Presentation
from sentence_transformers import SentenceTransformer, util
import difflib

###########################################
# 1. Text Extraction Functions for Various File Types
###########################################

def extract_txt(file_path):
    """Extracts text from a plain text (.txt) file."""
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

def extract_html(file_path):
    """Extracts text from an HTML file by stripping HTML tags."""
    with open(file_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f, 'html.parser')
    return soup.get_text(separator=' ', strip=True)

def extract_excel(file_path):
    """Extracts text from an Excel file (.xlsx or .xls) by reading its first sheet."""
    df = pd.read_excel(file_path)
    return df.to_string()

def extract_docx(file_path):
    """Extracts text from a Word document (.docx)."""
    doc = docx.Document(file_path)
    full_text = [para.text for para in doc.paragraphs]
    return "\n".join(full_text)

def extract_pdf(file_path):
    """Extracts text from a PDF file."""
    text = ""
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text

def extract_pptx(file_path):
    """Extracts text from a PowerPoint file (.pptx)."""
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

###########################################
# 2. Unified Text Extraction Based on File Extension
###########################################

def extract_text_from_file(file_path):
    """
    Determines the file type based on the extension and
    calls the corresponding extraction function.
    """
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.txt':
        return extract_txt(file_path)
    elif ext in ['.html', '.htm']:
        return extract_html(file_path)
    elif ext in ['.xlsx', '.xls']:
        return extract_excel(file_path)
    elif ext == '.docx':
        return extract_docx(file_path)
    elif ext == '.pdf':
        return extract_pdf(file_path)
    elif ext == '.pptx':
        return extract_pptx(file_path)
    else:
        raise ValueError("Unsupported file type")

###########################################
# 3. Comparison Functions Using AI/ML
###########################################

# Initialize the SentenceTransformer model once.
model = SentenceTransformer('all-MiniLM-L6-v2')

def compare_texts(text1, text2):
    """
    Compares two texts semantically using a pre-trained SentenceTransformer model.
    Returns a cosine similarity score between 0 and 1.
    """
    embedding1 = model.encode(text1, convert_to_tensor=True)
    embedding2 = model.encode(text2, convert_to_tensor=True)
    cosine_sim = util.cos_sim(embedding1, embedding2).item()
    return cosine_sim

def get_differences(text1, text2):
    """
    Generates a unified diff (line-by-line differences) between two texts.
    Returns the diff as a string.
    """
    diff = difflib.unified_diff(
        text1.splitlines(), text2.splitlines(),
        fromfile='File 1', tofile='File 2', lineterm=''
    )
    return "\n".join(list(diff))
def extract_txt(file_path):
    """Extracts text from a plain text (.txt) file using a fallback encoding if needed."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        # Fallback to a more lenient encoding; you can also try 'latin1' or 'cp1252'
        with open(file_path, 'r', encoding='latin1', errors='replace') as f:
            return f.read()
